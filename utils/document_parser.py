"""
Document Parser Module
Extracts text content from PDF and PPTX files.
"""

import fitz  # PyMuPDF
from pptx import Presentation
from typing import Dict
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def extract_from_pdf(file_path: str) -> Dict[str, any]:
    """
    Extract text from PDF slides.

    Args:
        file_path: Path to PDF file

    Returns:
        {
            'text': str,  # All extracted text
            'pages': int,  # Number of pages
            'has_images': bool  # Whether images were detected
        }
    """
    try:
        logger.info(f"Opening PDF file: {file_path}")
        doc = fitz.open(file_path)
        text_content = []
        has_images = False
        total_pages = len(doc)

        logger.info(f"PDF has {total_pages} pages")

        for page_num, page in enumerate(doc, 1):
            # Extract text from page
            page_text = page.get_text()
            logger.debug(f"Page {page_num} text length: {len(page_text)}")

            if page_text.strip():
                text_content.append(f"--- Page {page_num} ---\n{page_text}")

            # Check for images (optional for MVP)
            image_list = page.get_images()
            if image_list:
                has_images = True
                logger.debug(f"Page {page_num} has {len(image_list)} images")

        doc.close()

        total_text = '\n\n'.join(text_content)
        logger.info(f"Successfully extracted {len(total_text)} characters from {total_pages} pages")

        result = {
            'text': total_text,
            'pages': total_pages,
            'has_images': has_images
        }

        return result

    except Exception as e:
        logger.error(f"Error extracting PDF content: {str(e)}", exc_info=True)
        return {
            'text': '',
            'pages': 0,
            'has_images': False,
            'error': str(e)
        }


def extract_from_pptx(file_path: str) -> Dict[str, any]:
    """
    Extract text from PowerPoint slides.

    Args:
        file_path: Path to PPTX file

    Returns:
        {
            'text': str,  # All extracted text from slides
            'slides': int,  # Number of slides
            'notes': str  # Speaker notes if available
        }
    """
    try:
        prs = Presentation(file_path)
        text_content = []
        notes_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_content.append(f"--- Slide {slide_num} ---\n" + '\n'.join(slide_text))

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    notes_content.append(f"Notes for Slide {slide_num}: {notes_frame.text}")

        result = {
            'text': '\n\n'.join(text_content),
            'slides': len(prs.slides),
            'notes': '\n\n'.join(notes_content) if notes_content else ''
        }

        logger.info(f"Successfully extracted text from PPTX: {len(prs.slides)} slides")
        return result

    except Exception as e:
        logger.error(f"Error extracting PPTX content: {str(e)}")
        return {
            'text': '',
            'slides': 0,
            'notes': '',
            'error': str(e)
        }


def extract_from_transcript(file_path: str) -> Dict[str, any]:
    """
    Extract text from transcript files (TXT or VTT format).

    Args:
        file_path: Path to transcript file

    Returns:
        {
            'text': str,  # All extracted text
            'format': str,  # File format (txt or vtt)
            'lines': int  # Number of lines
        }
    """
    try:
        logger.info(f"Opening transcript file: {file_path}")

        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Determine if it's a VTT file
        file_ext = file_path.rsplit('.', 1)[1].lower() if '.' in file_path else 'txt'
        is_vtt = file_ext == 'vtt'

        if is_vtt:
            # Parse VTT format - remove timestamps and metadata
            lines = content.split('\n')
            text_lines = []
            skip_next = False

            for line in lines:
                line = line.strip()

                # Skip WEBVTT header
                if line.startswith('WEBVTT'):
                    continue

                # Skip timestamp lines (contain -->)
                if '-->' in line:
                    skip_next = False
                    continue

                # Skip cue identifiers (numbers or IDs before timestamps)
                if line and not skip_next and line[0].isdigit() and ':' not in line:
                    skip_next = True
                    continue

                # Add actual transcript text
                if line and not line.startswith('NOTE'):
                    text_lines.append(line)
                    skip_next = False

            cleaned_text = '\n'.join(text_lines)
        else:
            # Plain text file - use as-is
            cleaned_text = content

        line_count = len(cleaned_text.split('\n'))
        logger.info(f"Successfully extracted {len(cleaned_text)} characters from transcript ({line_count} lines)")

        result = {
            'text': cleaned_text,
            'format': file_ext,
            'lines': line_count
        }

        return result

    except Exception as e:
        logger.error(f"Error extracting transcript content: {str(e)}", exc_info=True)
        return {
            'text': '',
            'format': '',
            'lines': 0,
            'error': str(e)
        }


def extract_content(file_path: str, file_type: str) -> Dict[str, any]:
    """
    Main entry point - routes to appropriate extractor.

    Args:
        file_path: Path to uploaded file
        file_type: 'pdf', 'pptx', 'txt', or 'vtt'

    Returns:
        Extracted content dictionary
    """
    file_type = file_type.lower()

    if file_type == 'pdf':
        return extract_from_pdf(file_path)
    elif file_type in ['pptx', 'ppt']:
        return extract_from_pptx(file_path)
    elif file_type in ['txt', 'vtt']:
        return extract_from_transcript(file_path)
    else:
        logger.error(f"Unsupported file type: {file_type}")
        return {
            'text': '',
            'error': f'Unsupported file type: {file_type}'
        }
